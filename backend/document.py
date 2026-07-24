"""
document.py — extracts plain text from every file type the frontend
accepts, so extracted content can be prepended to a chat request.
Nothing here calls an LLM; it only turns files into text.
"""
from pathlib import Path

import openpyxl
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from pypdf.errors import PdfReadError

# Optional OCR dependencies (graceful degradation if not installed)
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:  # noqa: BLE001 — graceful degradation
    OCR_AVAILABLE = False

# Plain-text / source-code formats: read as-is, no library needed.
PLAIN_TEXT_EXTENSIONS = {
    "txt", "json", "html", "xml", "py", "java", "js", "c", "cpp",
    "cs", "go", "rs", "php", "sql", "r", "md",
}

# Image formats supported by Pillow + tesseract
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "tiff", "tif", "bmp", "webp", "gif"}

# Maximum number of pages to OCR in a scanned PDF (safety limit)
MAX_OCR_PAGES = 10

# Maximum file size to attempt OCR on (MB) — prevents OOM on huge scans
MAX_OCR_FILE_SIZE_MB = 50


def extract_text(path: Path, extension: str) -> str:
    """Dispatches to the right extractor for `extension`.
    Returns an empty string (rather than raising) for anything unexpected,
    so a bad file never takes down a chat request — the caller can decide
    how to warn the user."""
    extension = extension.lower().lstrip(".")

    try:
        if extension in PLAIN_TEXT_EXTENSIONS:
            return _extract_plain_text(path)
        if extension == "pdf":
            return _extract_pdf(path)
        if extension == "docx":
            return _extract_docx(path)
        if extension == "csv":
            return _extract_csv(path)
        if extension == "xlsx":
            return _extract_xlsx(path)
        if extension == "pptx":
            return _extract_pptx(path)
        if extension in IMAGE_EXTENSIONS:
            return _extract_image_ocr(path)
    except Exception as exc:  # noqa: BLE001 — surfaced to the caller as empty text + logged
        return f"[Could not extract text from this {extension} file: {exc}]"

    return ""


def _extract_plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF. Falls back to OCR for scanned/image-only PDFs."""
    try:
        reader = PdfReader(str(path))
        text_pages = [page.extract_text() or "" for page in reader.pages]
        extracted = "\n\n".join(text_pages)

        # If the extracted text is nearly empty, the PDF is likely scanned — try OCR
        if len(extracted.strip()) < 100 and OCR_AVAILABLE:
            ocr_text = _extract_pdf_ocr(path)
            if ocr_text and len(ocr_text.strip()) > len(extracted.strip()):
                return f"[OCR extracted from scanned PDF]\n\n{ocr_text}"

        return extracted
    except PdfReadError as exc:
        # Password-protected or corrupt PDF
        if "password" in str(exc).lower():
            return "[Password-protected PDF — please provide an unprotected copy]"
        # Try OCR as last resort for corrupt PDFs that might be scanned images
        if OCR_AVAILABLE:
            try:
                ocr_text = _extract_pdf_ocr(path)
                if ocr_text:
                    return f"[OCR extracted from PDF (could not parse normally): {exc}]\n\n{ocr_text}"
            except Exception:  # noqa: BLE001
                pass
        return f"[Could not read PDF: {exc}]"
    except Exception as exc:  # noqa: BLE001
        return f"[Could not extract text from this PDF file: {exc}]"


def _extract_pdf_ocr(path: Path) -> str:
    """Extract text from PDF using OCR (for scanned documents).
    Requires system tesseract binary and pytesseract + Pillow."""
    if not OCR_AVAILABLE:
        return ""

    # Safety check on file size
    size_mb = path.stat().st_size / (1024 * 1024)
    if size_mb > MAX_OCR_FILE_SIZE_MB:
        return f"[PDF too large for OCR ({size_mb:.1f} MB > {MAX_OCR_FILE_SIZE_MB} MB limit)]"

    # Try importing pdf2image (optional)
    try:
        from pdf2image import convert_from_path
    except ImportError:
        return "[PDF OCR requires pdf2image — install with 'pip install pdf2image']"

    # Convert limited number of pages to images for OCR
    images = convert_from_path(
        str(path),
        first_page=1,
        last_page=min(len(PdfReader(str(path)).pages), MAX_OCR_PAGES),
        dpi=200,
    )

    texts = []
    for i, image in enumerate(images, start=1):
        try:
            text = pytesseract.image_to_string(image)
            if text.strip():
                texts.append(f"--- Page {i} (OCR) ---\n{text}")
        except Exception:  # noqa: BLE001
            continue

    return "\n\n".join(texts)


def _extract_docx(path: Path) -> str:
    try:
        doc = DocxDocument(str(path))
        texts = [p.text for p in doc.paragraphs if p.text]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        texts.append(cell.text)
        return "\n".join(texts)
    except Exception as exc:  # noqa: BLE001
        # surfaced to caller
        return f"[Could not extract text from this DOCX file: {exc}]"


def _extract_csv(path: Path) -> str:
    try:
        df = pd.read_csv(path, nrows=200)
        return df.to_string(index=False, max_rows=200)
    except pd.errors.EmptyDataError:
        return "[CSV file is empty]"
    except pd.errors.ParserError as exc:
        return f"[Could not parse CSV: {exc}]"
    except Exception as exc:  # noqa: BLE001
        return f"[Could not extract text from this CSV file: {exc}]"


def _extract_xlsx(path: Path) -> str:
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        chunks = []
        for sheet in wb.worksheets:
            chunks.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True, max_row=200):
                chunks.append(", ".join("" if v is None else str(v) for v in row))
        return "\n".join(chunks)
    except Exception as exc:  # noqa: BLE001
        return f"[Could not extract text from this XLSX file: {exc}]"


def _extract_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(prs.slides, start=1):
            chunks.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text)
        return "\n".join(chunks)
    except Exception as exc:  # noqa: BLE001
        return f"[Could not extract text from this PPTX file: {exc}]"


def _extract_image_ocr(path: Path) -> str:
    """Extract text from an image file using OCR."""
    if not OCR_AVAILABLE:
        return "[OCR not available: install pytesseract and system tesseract to extract text from images]"

    try:
        # Safety check on file size
        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > MAX_OCR_FILE_SIZE_MB:
            return f"[Image too large for OCR ({size_mb:.1f} MB > {MAX_OCR_FILE_SIZE_MB} MB limit)]"

        image = Image.open(path)
        # Convert to RGB if needed (e.g., RGBA, P mode)
        if image.mode not in ("RGB", "L"):
            image = image.convert("RGB")

        text = pytesseract.image_to_string(image)
        if not text.strip():
            return "[No text detected in image]"
        return f"[OCR extracted from image]\n\n{text}"
    except Exception as exc:  # noqa: BLE001
        return f"[Could not extract text from this image: {exc}]"


def truncate_preview(text: str, length: int = 300) -> str:
    text = text.strip()
    return text if len(text) <= length else text[:length].rstrip() + "…"